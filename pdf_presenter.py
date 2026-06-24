import streamlit as st
import pandas as pd
import fitz  # pymupdf
import re
import io
import os
import psycopg2
from datetime import date
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

DB_HOST     = "db.maximaapparel.com"
DB_PORT     = 5432
DB_NAME     = "maxima_reporting"
DB_USER     = "agent_ro"
DB_PASSWORD = "R3@d1234!1"
TABLE       = "public.sss_upc_report"

CURRENCY_MAP = {
    "USD": ("wholesale_usd", "msrp_usd"),
    "CAD": ("wholesale_cad", "msrp_cad"),
    "GBP": ("wholesale_gbp", "msrp_gbp"),
    "EUR": ("wholesale_eur", "msrp_eur"),
    "AED": ("wholesale_aed", "msrp_aed"),
    "MXN": ("wholesale_mxn", "msrp_mxn"),
    "BRL": ("wholesale_brl", "msrp_brl"),
    "CLP": ("wholesale_clp", "msrp_clp"),
    "AUD": ("wholesale_aud", "msrp_aud"),
    "NZD": ("wholesale_nzd", "msrp_nzd"),
    "RMB": ("wholesale_rmb", "msrp_rmb"),
    "ARS": ("wholesale_ars", "msrp_ars"),
    "ECU": ("wholesale_ecu", "msrp_ecu"),
    "BOB": ("wholesale_bob", "msrp_bob"),
    "PEN": ("wholesale_pen", "msrp_pen"),
}

CURRENCY_SYMBOLS = {
    "USD": "$",    "CAD": "CA$",  "GBP": "GBP ", "EUR": "EUR ",
    "AED": "AED ", "MXN": "MX$", "BRL": "R$",   "CLP": "CLP$",
    "AUD": "AU$",  "NZD": "NZ$", "RMB": "RMB ", "ARS": "AR$",
    "ECU": "$",    "BOB": "Bs.", "PEN": "S/.",
}

PRICE_RE = re.compile(r'WS:\s*\$?([\d,]+\.?\d*)\s*/\s*MSRP:\s*\$?([\d,]+\.?\d*)', re.IGNORECASE)
STYLE_RE  = re.compile(r'^[A-Z]{2,5}\d{6,}-[A-Z0-9]{2,5}$')

st.set_page_config(page_title="PDF → PPT Converter", page_icon="📑", layout="centered")
st.markdown("""
<style>
@import url('https://fonts.googleapis.com/css2?family=Bebas+Neue&family=DM+Sans:wght@300;400;500&display=swap');
html,body,[class*="css"]{font-family:'DM Sans',sans-serif;}
.main-title{font-family:'Bebas Neue',sans-serif;font-size:3rem;letter-spacing:4px;color:#e8c84a;}
.sub-title{font-family:'Bebas Neue',sans-serif;font-size:1.1rem;letter-spacing:3px;color:#888;}
.success-box{background:rgba(74,222,128,0.1);border:1px solid rgba(74,222,128,0.4);border-radius:6px;padding:12px 16px;font-size:0.9rem;color:#4ade80;}
</style>
""", unsafe_allow_html=True)

st.markdown('<div class="main-title">PDF → PPT CONVERTER</div>', unsafe_allow_html=True)
st.markdown('<div class="sub-title">MAXIMA APPAREL — PRICE CONVERTER</div>', unsafe_allow_html=True)
st.markdown("---")


# ── Helpers ──────────────────────────────────────────────────────────────────
def fetch_prices(style_codes):
    if not style_codes:
        return {}
    price_cols = [col for pair in CURRENCY_MAP.values() for col in pair]
    try:
        conn = psycopg2.connect(
            host=DB_HOST, port=DB_PORT, dbname=DB_NAME,
            user=DB_USER, password=DB_PASSWORD, connect_timeout=15
        )
        cols_str = ", ".join(["ivnum"] + price_cols)
        df = pd.read_sql(
            f"SELECT DISTINCT {cols_str} FROM {TABLE} WHERE ivnum = ANY(%s)",
            conn, params=(list(style_codes),)
        )
        conn.close()
        return {row["ivnum"]: row for _, row in df.iterrows()}
    except Exception as e:
        st.error(f"Error DB: {e}")
        return {}


def int_to_rgb(c):
    return ((c >> 16) & 0xFF) / 255, ((c >> 8) & 0xFF) / 255, (c & 0xFF) / 255


def fmt_price(val, currency):
    sym = CURRENCY_SYMBOLS.get(currency, "$")
    if val is None or (isinstance(val, float) and pd.isna(val)):
        return "N/A"
    if currency == "CLP":
        return f"{sym}{int(float(val)):,}"
    return f"{sym}{float(val):,.2f}"


def extract_styles_from_pdf(pdf_bytes):
    """Returns list of (style_code) found in PDF, in order."""
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    styles = []
    seen = set()
    for pn in range(len(doc)):
        for line in doc[pn].get_text().splitlines():
            s = line.strip()
            if STYLE_RE.match(s) and s not in seen:
                styles.append(s)
                seen.add(s)
    doc.close()
    return styles


def process_pdf(pdf_bytes, price_data, currency):
    """Replace prices in PDF. price_data: {ivnum: {ws_col: val, ms_col: val}}"""
    ws_col, ms_col = CURRENCY_MAP[currency]
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    result_pngs = []

    for page_num in range(len(doc)):
        page = doc[page_num]
        if not page.get_text().strip():
            continue

        # Collect spans
        spans = []
        for block in page.get_text("dict", flags=fitz.TEXT_PRESERVE_WHITESPACE)["blocks"]:
            if block.get("type") != 0:
                continue
            for line in block.get("lines", []):
                for span in line.get("spans", []):
                    if span["text"].strip():
                        spans.append(span)

        replacements = []
        current_style = None
        for span in spans:
            txt = span["text"].strip()
            if STYLE_RE.match(txt):
                current_style = txt
            elif PRICE_RE.search(txt) and current_style:
                row = price_data.get(current_style)
                if row is not None:
                    new_text = (
                        f"WS: {fmt_price(row.get(ws_col), currency)} / "
                        f"MSRP: {fmt_price(row.get(ms_col), currency)}"
                    )
                    replacements.append({
                        "rect":     fitz.Rect(span["bbox"]),
                        "new_text": new_text,
                        "color":    int_to_rgb(span["color"]),
                        "size":     span["size"],
                    })
                current_style = None

        # Redact with WHITE background (price areas are always on white)
        for rep in replacements:
            page.add_redact_annot(rep["rect"], fill=(1, 1, 1))
        page.apply_redactions()

        # Re-insert text
        for rep in replacements:
            rect = rep["rect"]
            page.insert_text(
                (rect.x0, rect.y1 - 1),
                rep["new_text"],
                fontsize=rep["size"],
                color=rep["color"],
                fontname="helv",
            )

        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
        result_pngs.append(pix.tobytes("png"))

    doc.close()
    return result_pngs


def build_pptx(pngs_by_currency):
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for currency, pngs in pngs_by_currency.items():
        for png_bytes in pngs:
            slide = prs.slides.add_slide(blank)
            img = Image.open(io.BytesIO(png_bytes))
            w, h = img.size
            aspect = w / h
            sw, sh = prs.slide_width, prs.slide_height
            if aspect > sw / sh:
                pw, ph = sw, int(sw / aspect)
            else:
                ph, pw = sh, int(sh * aspect)
            left = (sw - pw) // 2
            top  = (sh - ph) // 2
            slide.shapes.add_picture(io.BytesIO(png_bytes), left, top, pw, ph)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ── UI ────────────────────────────────────────────────────────────────────────
st.markdown("#### Paso 1 — Sube los PDFs")
uploaded_pdfs = st.file_uploader(
    "Arrastra o selecciona los PDFs de presentación",
    type=["pdf"], accept_multiple_files=True,
)

st.markdown("#### Paso 2 — Selecciona moneda(s)")
selected_currencies = st.multiselect(
    "Moneda(s)", options=list(CURRENCY_MAP.keys()), default=["USD"],
    placeholder="Selecciona moneda(s)...",
)

if uploaded_pdfs and selected_currencies:
    # Read PDFs and extract styles
    pdf_store = {}
    all_styles_ordered = []
    seen_styles = set()
    for f in uploaded_pdfs:
        pdf_bytes = f.read()
        styles = extract_styles_from_pdf(pdf_bytes)
        pdf_store[f.name] = pdf_bytes
        for s in styles:
            if s not in seen_styles:
                all_styles_ordered.append(s)
                seen_styles.add(s)

    st.markdown("---")
    tab_auto, tab_manual = st.tabs(["🔄 Automático (Base de datos)", "✏️ Manual (ingresar precios)"])

    # ── TAB AUTOMÁTICO ────────────────────────────────────────────────────────
    with tab_auto:
        st.markdown(f"**{len(all_styles_ordered)} estilo(s) detectados:** {', '.join(all_styles_ordered)}")

        if st.button("🔄 CONVERTIR Y GENERAR PPT", type="primary", use_container_width=True, key="btn_auto"):
            st.session_state.pop("ppt_auto", None)

            with st.spinner("Consultando base de datos..."):
                price_data = fetch_prices(set(all_styles_ordered))

            found   = set(price_data.keys())
            missing = [s for s in all_styles_ordered if s not in found]

            st.info(f"✅ {len(found)} encontrados en DB" + (
                f" | ⚠️ No encontrados (usa tab Manual): {', '.join(missing)}" if missing else ""
            ))

            pngs_by_currency = {c: [] for c in selected_currencies}
            for pdf_name, pdf_bytes in pdf_store.items():
                for currency in selected_currencies:
                    with st.spinner(f"{pdf_name} → {currency}..."):
                        pngs = process_pdf(pdf_bytes, price_data, currency)
                        pngs_by_currency[currency].extend(pngs)

            total = sum(len(v) for v in pngs_by_currency.values())
            with st.spinner("Generando PPT..."):
                pptx_bytes = build_pptx(pngs_by_currency)

            curr_label = "_".join(selected_currencies)
            fname = f"Presentacion_{curr_label}_{date.today().strftime('%m.%d.%Y')}.pptx"
            out_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "UPCs_generados")
            os.makedirs(out_dir, exist_ok=True)
            try:
                with open(os.path.join(out_dir, fname), "wb") as fh:
                    fh.write(pptx_bytes)
            except Exception:
                pass
            st.session_state["ppt_auto"] = {"fname": fname, "bytes": pptx_bytes, "slides": total}

        out = st.session_state.get("ppt_auto")
        if out:
            st.markdown(f'<div class="success-box">✅ {out["slides"]} diapositiva(s) generada(s)</div>', unsafe_allow_html=True)
            st.download_button(
                label=f"⬇ Descargar {out['fname']}",
                data=out["bytes"], file_name=out["fname"],
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True, key="dl_auto",
            )

    # ── TAB MANUAL ────────────────────────────────────────────────────────────
    with tab_manual:
        st.markdown("Sube tu archivo Excel con los precios. Columnas requeridas: **STYLE** o **STOCK**, **CURRENCY**, **LINE COST**, **MSRP**.")

        price_file = st.file_uploader("Sube el Excel con precios", type=["xlsx", "xls", "csv"], key="price_file")

        if price_file:
            try:
                if price_file.name.endswith(".csv"):
                    df_prices = pd.read_csv(price_file, dtype=str)
                else:
                    df_prices = pd.read_excel(price_file, dtype=str)

                df_prices.columns = [str(c).strip().upper() for c in df_prices.columns]

                style_col_xls = next((c for c in ["STOCK", "STYLE"] if c in df_prices.columns), None)
                currency_col  = next((c for c in ["CURRENCY", "MONEDA"] if c in df_prices.columns), None)
                ws_col_xls    = next((c for c in ["LINE COST", "LINE_COST", "WS", "WHOLESALE"] if c in df_prices.columns), None)
                ms_col_xls    = next((c for c in ["MSRP"] if c in df_prices.columns), None)

                if not all([style_col_xls, ws_col_xls, ms_col_xls]):
                    st.error(f"Columnas no encontradas. Detectadas: {', '.join(df_prices.columns)}")
                else:
                    st.success(f"✅ {len(df_prices)} filas cargadas — columna estilo: **{style_col_xls}**")

                    if currency_col:
                        currencies_in_file = df_prices[currency_col].dropna().unique().tolist()
                        file_currency = currencies_in_file[0].strip().upper() if currencies_in_file else "USD"
                    else:
                        file_currency = "USD"

                    st.info(f"Moneda detectada: **{file_currency}**")

                    ws_db, ms_db = CURRENCY_MAP.get(file_currency, ("wholesale_usd", "msrp_usd"))
                    price_data_manual = {}
                    for _, row in df_prices.iterrows():
                        style_key = str(row.get(style_col_xls, "") or "").strip()
                        if not style_key:
                            continue
                        try:
                            ws_val = float(str(row.get(ws_col_xls, "") or "").replace(",", ""))
                        except (ValueError, TypeError):
                            ws_val = None
                        try:
                            ms_val = float(str(row.get(ms_col_xls, "") or "").replace(",", ""))
                        except (ValueError, TypeError):
                            ms_val = None
                        price_data_manual[style_key] = {ws_db: ws_val, ms_db: ms_val}

                    matched   = [s for s in all_styles_ordered if s in price_data_manual]
                    unmatched = [s for s in all_styles_ordered if s not in price_data_manual]
                    st.info(f"✅ {len(matched)} estilo(s) coinciden con el PDF" + (
                        f" | ⚠️ Sin precio: {', '.join(unmatched)}" if unmatched else ""
                    ))

                    if st.button("📄 GENERAR PPT CON ARCHIVO DE PRECIOS", type="primary", use_container_width=True, key="btn_manual"):
                        st.session_state.pop("ppt_manual", None)
                        pngs_by_currency = {file_currency: []}
                        for pdf_name, pdf_bytes in pdf_store.items():
                            with st.spinner(f"Procesando {pdf_name}..."):
                                pngs = process_pdf(pdf_bytes, price_data_manual, file_currency)
                                pngs_by_currency[file_currency].extend(pngs)
                        total = sum(len(v) for v in pngs_by_currency.values())
                        with st.spinner("Generando PPT..."):
                            pptx_bytes = build_pptx(pngs_by_currency)
                        fname = f"Presentacion_{file_currency}_{date.today().strftime('%m.%d.%Y')}.pptx"
                        out_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "UPCs_generados")
                        os.makedirs(out_dir, exist_ok=True)
                        try:
                            with open(os.path.join(out_dir, fname), "wb") as fh:
                                fh.write(pptx_bytes)
                        except Exception:
                            pass
                        st.session_state["ppt_manual"] = {"fname": fname, "bytes": pptx_bytes, "slides": total}

            except Exception as e:
                st.error(f"Error leyendo el archivo: {e}")

        out_m = st.session_state.get("ppt_manual")
        if out_m:
            st.markdown(f'<div class="success-box">✅ {out_m["slides"]} diapositiva(s) generada(s)</div>', unsafe_allow_html=True)
            st.download_button(
                label=f"⬇ Descargar {out_m['fname']}",
                data=out_m["bytes"], file_name=out_m["fname"],
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    use_container_width=True, key="dl_manual",
                )
