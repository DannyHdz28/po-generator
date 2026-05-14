import io
import re
from collections import defaultdict
from pathlib import Path

import streamlit as st
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from pptx.util import Emu

# ─── PAGE CONFIG ──────────────────────────────────────────────
st.set_page_config(page_title="Deck Builder — Pro Standard", page_icon="🏆", layout="wide")

st.markdown("""
<style>
@import url('https://fonts.googleapis.com/css2?family=Bebas+Neue&family=DM+Sans:wght@300;400;500&display=swap');
html, body, [class*="css"] { font-family: 'DM Sans', sans-serif; }
.main-title  { font-family: 'Bebas Neue', sans-serif; font-size: 3rem;   letter-spacing: 4px; color: #e8c84a; margin-bottom: 0; }
.sub-title   { font-family: 'Bebas Neue', sans-serif; font-size: 1.1rem; letter-spacing: 3px; color: #888;     margin-top: 0; }
.section-h   { font-family: 'Bebas Neue', sans-serif; font-size: 1.4rem; letter-spacing: 2px; color: #f5f5f0;  margin-top: 1.5rem; }
.pill-ok     { background: rgba(74,222,128,0.12); border:1px solid rgba(74,222,128,0.4); color:#4ade80; padding:2px 10px; border-radius:3px; font-size:0.75rem; letter-spacing:1px; }
.pill-warn   { background: rgba(248,113,113,0.12); border:1px solid rgba(248,113,113,0.4); color:#f87171; padding:2px 10px; border-radius:3px; font-size:0.75rem; letter-spacing:1px; }
.capsule-card { background:#1a1a1a; border:1px solid #2a2a2a; border-radius:4px; padding:10px 14px; }
.capsule-card.match { border-color:#4ade80; }
.capsule-card h4 { font-family:'Bebas Neue',sans-serif; font-size:1rem; letter-spacing:1.5px; color:#f5f5f0; margin:0; }
.capsule-card p  { font-size:0.75rem; color:#888; margin:2px 0 0 0; }
.team-badge { display:inline-block; padding:6px 14px; background:rgba(232,200,74,0.1);
              border:1px solid rgba(232,200,74,0.3); border-radius:3px; font-size:0.85rem; }
</style>
""", unsafe_allow_html=True)

# ─── CONSTANTS ────────────────────────────────────────────────
GENDER_CODES       = {"M", "W", "K"}
IMAGE_EXTS         = {".jpg", ".jpeg", ".png"}
CATEGORY_TO_GENDER = {"MENS": "M", "WOMENS": "W", "KIDS": "K", "BOYS": "K", "GIRLS": "K"}

# Ruta base por defecto del servidor de la empresa
SERVER_BASE_DEFAULT = r"\\10.0.1.30\Sales Toolkits\PROSTANDARD\USA_CANADA\CATALOGS"

# Cápsulas por cuarto. Q2-Q4 se llenan cuando lleguen los nombres.
CAPSULES_BY_QUARTER = {
    "Q1": [
        "TEAM CITY", "FLAGSHIP", "SPRING BREAK", "ULTIMATE FAN",
        "PRO FILE", "PROPERTY OF", "HERITAGE HUSTLE",
        "HERITAGE & HUSTLE", "REFLECTION", "FLORAL SPORT",
        "CLASSIC ICON", "CLASSICS",
    ],
    "Q2": [],
    "Q3": [],
    "Q4": [],
}

# Tamaño estándar para slides sin imagen previa (29.44cm × 19.05cm @ x=2.21cm, y=0)
INSERT_LEFT   = Emu(795600)
INSERT_TOP    = Emu(0)
INSERT_WIDTH  = Emu(10598400)
INSERT_HEIGHT = Emu(6858000)


# ─── PARSERS ──────────────────────────────────────────────────
def norm(s: str) -> str:
    return re.sub(r"\s+", " ", re.sub(r"[&\-]", " ", s.upper())).strip()


def parse_merch_name(filename: str):
    """Parsea LIGA_EQUIPO_CAPSULA[_GENERO][_NNN].ext"""
    base = re.sub(r"\.[^.]+$", "", filename)
    parts = base.split("_")
    if len(parts) < 3:
        return None
    liga    = parts[0].upper().strip()
    end_idx = len(parts) - 1
    if parts[end_idx].isdigit():
        end_idx -= 1
    gender = None
    if end_idx >= 0 and parts[end_idx].upper() in GENDER_CODES:
        gender  = parts[end_idx].upper()
        end_idx -= 1
    if end_idx < 1:
        return None
    team    = parts[1].strip()
    capsule = " ".join(parts[2:end_idx + 1]).upper().strip()
    if not capsule:
        return None
    key = capsule + (f" {gender}" if gender else "")
    return {"liga": liga, "team": team, "capsule": capsule, "gender": gender, "key": key}


def parse_note_line(line: str):
    """Parsea una línea de notas del PPT → {key, capsule, gender}"""
    n = norm(line).strip()
    n = re.sub(r"\s+\d+$", "", n).strip()
    if len(n) < 2:
        return None
    parts = n.split(" ")
    last  = parts[-1]
    if last in GENDER_CODES:
        return {"key": n, "capsule": " ".join(parts[:-1]), "gender": last}
    return {"key": n, "capsule": n, "gender": None}


CAT_FOLDER_PREFIXES = {
    "MENS":   ["MENS", "MNS"],
    "WOMENS": ["WOMENS", "WMNS", "WNS"],
    "KIDS":   ["KIDS", "KDS"],
    "BOYS":   ["BOYS"],
    "GIRLS":  ["GIRLS"],
}


def extract_capsule_from_folder(folder_name: str, category: str) -> str:
    """Extrae clave de cápsula del nombre de carpeta del servidor.
    Ej: '01·01 - Q1 2027 MENS TEAM CITY' + 'MENS' → 'TEAM CITY M'
         '01·01 - Q1 2027 WMNS FLORAL SPORT' + 'WOMENS' → 'FLORAL SPORT W'
    """
    n = re.sub(r'^.+?[-–]\s*', '', folder_name).strip()
    n = re.sub(r'^Q\d\s+\d{4}\s+', '', n, flags=re.IGNORECASE).strip()
    cat_upper = category.upper()
    for prefix in CAT_FOLDER_PREFIXES.get(cat_upper, [cat_upper]):
        if n.upper().startswith(prefix + " "):
            n = n[len(prefix):].strip()
            break
        elif n.upper() == prefix:
            n = ""
            break
    gender = CATEGORY_TO_GENDER.get(cat_upper)
    key = n.upper()
    if gender and key:
        key = f"{key} {gender}"
    return key or folder_name.upper()


def find_note_match(notes_text: str, merch_map: dict):
    """Primera línea de notas que matchee una key del merch_map."""
    if not notes_text:
        return None
    for raw in notes_text.split("\n"):
        p = parse_note_line(raw)
        if not p:
            continue
        # 1. Direct key match
        if p["key"] in merch_map:
            return p["key"]
        # 2. Capsule match
        if p["capsule"] in merch_map:
            return p["capsule"]
        # 3. Normalized full-key match
        norm_p = norm(p["key"])
        for k in merch_map:
            if norm(k) == norm_p:
                return k
        # 4. Gender-agnostic: solo cuando la nota NO tiene género especificado
        if not p["gender"]:
            norm_capsule = norm(p["capsule"])
            for k in merch_map:
                k_parts = norm(k).split()
                k_cap = " ".join(k_parts[:-1]) if k_parts and k_parts[-1] in GENDER_CODES else norm(k)
                if k_cap == norm_capsule:
                    return k
        # 5. Partial word match: nota abreviada (ej: "HERITAGE M" → "HERITAGE & HUSTLE M")
        note_cap_words = norm(p["capsule"]).split()
        for k in merch_map:
            k_parts = norm(k).split()
            k_gender = k_parts[-1] if k_parts and k_parts[-1] in GENDER_CODES else None
            k_cap_words = k_parts[:-1] if k_gender else k_parts
            if (k_cap_words[:len(note_cap_words)] == note_cap_words
                    and (not p["gender"] or p["gender"] == k_gender)):
                return k
    return None


# ─── SERVER HELPERS ───────────────────────────────────────────
def find_quarter_dir(year_path: Path, quarter: str) -> Path | None:
    """Encuentra la carpeta del cuarto, maneja 'Q1' y 'Q1 2027'."""
    direct = year_path / quarter
    if direct.is_dir():
        return direct
    try:
        for d in year_path.iterdir():
            if d.is_dir() and d.name.upper().startswith(quarter.upper()):
                return d
    except Exception:
        pass
    return None


def list_dirs(path: Path) -> list[str]:
    """Lista nombres de subcarpetas ordenados."""
    try:
        return sorted([d.name for d in path.iterdir() if d.is_dir()])
    except PermissionError:
        st.error(f"Sin permisos para leer: {path}")
        return []
    except Exception as e:
        st.error(f"Error leyendo carpeta: {e}")
        return []


def list_images(path: Path) -> list[Path]:
    """Lista archivos de imagen ordenados por nombre."""
    try:
        return sorted(
            [f for f in path.iterdir() if f.is_file() and f.suffix.lower() in IMAGE_EXTS]
        )
    except Exception:
        return []


def find_merchboards_dir(capsule_path: Path) -> Path | None:
    """Busca la carpeta _MERCHBOARDS dentro de la cápsula."""
    # intento directo
    direct = capsule_path / "_MERCHBOARDS"
    if direct.is_dir():
        return direct
    # búsqueda flexible (por si el nombre varía un poco)
    try:
        for d in capsule_path.iterdir():
            if d.is_dir() and "MERCHBOARD" in d.name.upper():
                return d
    except Exception:
        pass
    return None


def build_merch_map_from_paths(image_paths: list[Path]) -> tuple[dict, str | None, list[str]]:
    """Construye merch_map leyendo los archivos del servidor."""
    merch_map    = {}
    detected_team = None
    unmatched    = []

    for img_path in image_paths:
        parsed = parse_merch_name(img_path.name)
        if parsed and parsed["capsule"]:
            merch_map.setdefault(parsed["key"], []).append({
                "name":  img_path.name,
                "bytes": img_path.read_bytes(),
            })
            if not detected_team:
                detected_team = parsed["team"]
        else:
            unmatched.append(img_path.name)

    for k in merch_map:
        merch_map[k].sort(key=lambda x: x["name"])

    return merch_map, detected_team, unmatched


# ─── PPT OPERATIONS ───────────────────────────────────────────
def scan_pptx(file_bytes: bytes, known_capsules: list | None = None):
    prs        = Presentation(io.BytesIO(file_bytes))
    known_norm = {norm(c) for c in (known_capsules or [])}
    found, missing_gender, unknown_capsule = [], [], []

    for idx, slide in enumerate(prs.slides, start=1):
        if not slide.has_notes_slide:
            continue
        notes = slide.notes_slide.notes_text_frame.text or ""
        for raw in notes.split("\n"):
            p = parse_note_line(raw)
            if not (p and p["capsule"]):
                continue
            found.append({"slide": idx, "key": p["key"], "capsule": p["capsule"], "gender": p["gender"]})
            if not p["gender"]:
                missing_gender.append({"slide": idx, "capsule": p["capsule"]})
            if known_norm and norm(p["capsule"]) not in known_norm:
                unknown_capsule.append({"slide": idx, "capsule": p["capsule"]})
            break

    return found, missing_gender, unknown_capsule


def replace_picture_blob(shape, new_blob: bytes) -> bool:
    blip = shape._element.find(".//" + qn("a:blip"))
    if blip is None:
        return False
    r_id = blip.get(qn("r:embed"))
    if not r_id:
        return False
    img_part       = shape.part._rels[r_id].target_part
    img_part._blob = new_blob

    # Remove crop (srcRect) so the full image shows
    blip_fill = blip.getparent()
    src_rect  = blip_fill.find(qn("a:srcRect"))
    if src_rect is not None:
        blip_fill.remove(src_rect)

    # Fix position and size to standard merchboard dimensions
    shape.left   = INSERT_LEFT
    shape.top    = INSERT_TOP
    shape.width  = INSERT_WIDTH
    shape.height = INSERT_HEIGHT

    return True


def generate_deck(ppt_bytes: bytes, merch_map: dict):
    prs      = Presentation(io.BytesIO(ppt_bytes))
    used     = defaultdict(int)
    replaced = 0
    log      = []

    for idx, slide in enumerate(prs.slides, start=1):
        if not slide.has_notes_slide:
            continue
        notes       = slide.notes_slide.notes_text_frame.text or ""
        matched_key = find_note_match(notes, merch_map)
        if not matched_key:
            continue

        imgs = merch_map[matched_key]
        if used[matched_key] >= len(imgs):
            log.append(("warn", f"Slide {idx} ({matched_key}): sin más imágenes disponibles"))
            continue

        img = imgs[used[matched_key]]
        used[matched_key] += 1

        pics = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]

        if pics:
            if replace_picture_blob(pics[0], img["bytes"]):
                replaced += 1
                log.append(("ok", f"Slide {idx} → {matched_key} → {img['name']}"))
            else:
                log.append(("err", f"Slide {idx} ({matched_key}): no se pudo reemplazar"))
        else:
            try:
                slide.shapes.add_picture(
                    io.BytesIO(img["bytes"]),
                    left=INSERT_LEFT, top=INSERT_TOP,
                    width=INSERT_WIDTH, height=INSERT_HEIGHT,
                )
                replaced += 1
                log.append(("ok", f"Slide {idx} → {matched_key} → {img['name']} (insertada)"))
            except Exception as e:
                log.append(("err", f"Slide {idx} ({matched_key}): error insertando — {e}"))

    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    return out, replaced, log


# ═══════════════════════════════════════════════════════════════
# UI
# ═══════════════════════════════════════════════════════════════
st.markdown('<div class="main-title">DECK BUILDER</div>', unsafe_allow_html=True)
st.markdown('<div class="sub-title">PRO STANDARD · MERCHBOARD AUTO-REPLACE</div>', unsafe_allow_html=True)
st.write("")

with st.expander("Cómo preparar tu PPT", expanded=False):
    st.markdown("""
1. En cada slide de **merchboard** agrega en las **Notas** la cápsula + género:
   `FLAGSHIP M`, `SPRING BREAK W`, `TEAM CITY K`  (M = Mens, W = Womens, K = Kids).
2. Las imágenes en el servidor deben seguir el naming:
   `LIGA_EQUIPO_CAPSULA_GENERO[_NNN].jpg`
   Ej: `NFL_DALLAS COWBOYS_TEAM CITY_M_001.jpg`
    """)

# ─── SIDEBAR ──────────────────────────────────────────────────
with st.sidebar:
    st.markdown("### CONFIG")

    quarter = st.selectbox("Cuarto", ["Q1", "Q2", "Q3", "Q4"], index=0)
    known_capsules = CAPSULES_BY_QUARTER.get(quarter, [])
    if known_capsules:
        st.caption(f"Cápsulas {quarter}:")
        for c in known_capsules:
            st.caption(f"· {c}")
    else:
        st.caption(f"Sin cápsulas definidas para {quarter}.")

    st.markdown("---")
    st.markdown("### SERVIDOR")
    server_base = st.text_input(
        "Ruta base",
        value=SERVER_BASE_DEFAULT,
        help="Ruta UNC al servidor, ej: \\\\10.0.1.30\\Sales Toolkits\\...",
    )

# ─── STEP 1: PPT BASE ─────────────────────────────────────────
st.markdown('<div class="section-h">1 · PPT BASE</div>', unsafe_allow_html=True)
ppt_file = st.file_uploader("Sube tu PPT base (.pptx)", type=["pptx"], key="ppt")

ppt_bytes    = None
slides_found = []

if ppt_file:
    ppt_bytes = ppt_file.getvalue()
    try:
        slides_found, missing_gender, unknown_capsule = scan_pptx(ppt_bytes, known_capsules)
        if slides_found:
            st.markdown(f'<span class="pill-ok">✓ {len(slides_found)} SLIDE(S) MARCADAS</span>', unsafe_allow_html=True)
            for s in slides_found:
                st.caption(f"Slide {s['slide']} → {s['capsule']} · género: {s['gender'] or '—'}")
        else:
            st.markdown('<span class="pill-warn">⚠ SIN NOTAS DE CÁPSULA</span>', unsafe_allow_html=True)
            st.caption("Verifica que las slides de merchboard tengan notas tipo `FLAGSHIP M`.")

        if missing_gender:
            st.warning(
                f"⚠ {len(missing_gender)} slide(s) sin género (M/W/K) en la nota:\n\n"
                + "\n".join(f"- Slide {m['slide']} → `{m['capsule']}`" for m in missing_gender)
            )
        if unknown_capsule:
            st.warning(
                f"⚠ {len(unknown_capsule)} slide(s) con cápsulas no listadas en {quarter}:\n\n"
                + "\n".join(f"- Slide {u['slide']} → `{u['capsule']}`" for u in unknown_capsule)
            )
    except Exception as e:
        st.error(f"Error leyendo el PPT: {e}")

# ─── STEP 2: MERCHBOARDS ──────────────────────────────────────
st.markdown('<div class="section-h">2 · MERCHBOARDS</div>', unsafe_allow_html=True)

modo = st.radio(
    "Origen de las imágenes",
    ["🖥️  Desde el servidor", "📁  Subir manualmente"],
    horizontal=True,
)

merch_map     = {}
detected_team = None
unmatched_names = []

# ── MODO SERVIDOR ─────────────────────────────────────────────
if modo == "🖥️  Desde el servidor":
    base_path = Path(server_base)

    if not base_path.exists():
        st.error(f"No se puede acceder a la ruta: `{server_base}`\n\nVerifica que estés conectado a la red de la empresa.")
    else:
        # Nivel 1: Categoría (multiselect)
        categories = list_dirs(base_path)
        if categories:
            sel_categories = st.multiselect(
                "Categoría (puedes seleccionar varias)", categories,
                default=categories[:1], key="srv_category",
                help="MENS, WOMENS, KIDS, etc."
            )
        else:
            st.warning("No se encontraron carpetas de categoría.")
            sel_categories = []

        # Nivel 2: Año (toma del primer cat que no sea KIDS, o del primer sub-cat de KIDS)
        sel_year = None
        if sel_categories:
            non_kids = [c for c in sel_categories if c.upper() != "KIDS"]
            if non_kids:
                year_base = base_path / non_kids[0]
            else:
                kids_subs = list_dirs(base_path / "KIDS")
                year_base = base_path / "KIDS" / kids_subs[0] if kids_subs else None
            if year_base and year_base.is_dir():
                years = list_dirs(year_base)
                if years:
                    sel_year = st.selectbox("Año", years, key="srv_year")

        # Nivel 2.5: Sub-categoría KIDS (solo si KIDS está seleccionado)
        kids_sub_selection = []
        if sel_year and any(c.upper() == "KIDS" for c in sel_categories):
            kids_subs = list_dirs(base_path / "KIDS")
            if kids_subs:
                kids_sub_selection = st.multiselect(
                    "Sub-categoría KIDS", kids_subs,
                    default=[], key="srv_kids_sub",
                    help="BOYS, GIRLS, KIDS"
                )

        # Construir fuentes efectivas: label → q_path
        sources = {}  # {label: quarter_path}
        if sel_year:
            for cat in sel_categories:
                if cat.upper() == "KIDS":
                    for sub in kids_sub_selection:
                        q_dir = find_quarter_dir(base_path / cat / sub / sel_year, quarter)
                        if q_dir:
                            sources[sub] = q_dir
                else:
                    q_dir = find_quarter_dir(base_path / cat / sel_year, quarter)
                    if q_dir:
                        sources[cat] = q_dir

        # Nivel 3: Cápsulas (multiselect — agrega carpetas de todas las fuentes)
        sel_capsule_folders = []
        capsule_folders_by_cat = {}  # label → {"q_path": Path, "folders": [...]}

        if sources:
            for label, q_path in sources.items():
                folders = list_dirs(q_path)
                if folders:
                    capsule_folders_by_cat[label] = {"q_path": q_path, "folders": folders}

            all_capsule_folders = sorted({
                f for info in capsule_folders_by_cat.values() for f in info["folders"]
            })

            if all_capsule_folders:
                sel_capsule_folders = st.multiselect(
                    "Cápsula(s) (puedes seleccionar varias)", all_capsule_folders,
                    default=[], key="srv_capsule",
                    help="Selecciona las cápsulas que quieres incluir en el deck"
                )
            else:
                st.warning(f"No hay carpetas de cápsula para {quarter} en las categorías seleccionadas.")

        # Nivel 4: Liga y Equipo
        sel_league = None
        sel_team   = None

        if sel_capsule_folders:
            leagues_set = set()
            for cap in sel_capsule_folders:
                for label, info in capsule_folders_by_cat.items():
                    if cap in info["folders"]:
                        mb = find_merchboards_dir(info["q_path"] / cap)
                        if mb:
                            for d in list_dirs(mb):
                                leagues_set.add(d)
            leagues = sorted(leagues_set)

            if leagues:
                sel_league = st.selectbox("Liga", leagues, key="srv_league")
            else:
                st.warning("No se encontró carpeta _MERCHBOARDS en las cápsulas seleccionadas.")

            if sel_league:
                teams_set = set()
                for cap in sel_capsule_folders:
                    for label, info in capsule_folders_by_cat.items():
                        if cap in info["folders"]:
                            mb = find_merchboards_dir(info["q_path"] / cap)
                            if mb:
                                league_path = mb / sel_league
                                if league_path.is_dir():
                                    for t in list_dirs(league_path):
                                        teams_set.add(t)
                teams = sorted(teams_set)

                if teams:
                    sel_team = st.selectbox("Equipo", teams, key="srv_team")
                else:
                    st.warning(f"No hay equipos en {sel_league}.")

        # Nivel 5: Cargar imágenes
        if sel_team and sel_league:
            loaded_routes = []

            for cap_folder in sel_capsule_folders:
                for label, info in capsule_folders_by_cat.items():
                    if cap_folder not in info["folders"]:
                        continue
                    mb = find_merchboards_dir(info["q_path"] / cap_folder)
                    if not mb:
                        continue
                    team_path = mb / sel_league / sel_team
                    images    = list_images(team_path)
                    if not images:
                        continue
                    cap_key = extract_capsule_from_folder(cap_folder, label)
                    loaded_routes.append(f"{label} / {cap_folder}")
                    for img in images:
                        merch_map.setdefault(cap_key, []).append({
                            "name":  img.name,
                            "bytes": img.read_bytes(),
                        })
                    detected_team = sel_team

            for k in merch_map:
                merch_map[k].sort(key=lambda x: x["name"])

            if detected_team:
                st.markdown(
                    f'<div class="team-badge">'
                    f'<span style="color:#888;font-size:0.7rem;letter-spacing:1px;">EQUIPO · </span>'
                    f'<span style="color:#e8c84a;font-family:\'Bebas Neue\',sans-serif;letter-spacing:1px;">'
                    f'{sel_team}</span></div>',
                    unsafe_allow_html=True,
                )
                if loaded_routes:
                    st.caption("Rutas cargadas: " + ", ".join(loaded_routes))
                st.write("")

            if merch_map:
                cols = st.columns(min(4, len(merch_map)))
                for i, (k, v) in enumerate(merch_map.items()):
                    with cols[i % len(cols)]:
                        st.markdown(
                            f'<div class="capsule-card match"><h4>{k}</h4>'
                            f'<p>{len(v)} imagen(es)</p></div>',
                            unsafe_allow_html=True,
                        )
            elif sel_team:
                st.warning(f"No se encontraron imágenes para {sel_team} en las cápsulas seleccionadas.")

# ── MODO MANUAL ───────────────────────────────────────────────
else:
    merch_files = st.file_uploader(
        "Sube las imágenes (JPG/PNG)",
        type=["jpg", "jpeg", "png"],
        accept_multiple_files=True,
        key="merch_manual",
    )

    if merch_files:
        for f in merch_files:
            parsed = parse_merch_name(f.name)
            if parsed and parsed["capsule"]:
                merch_map.setdefault(parsed["key"], []).append({"name": f.name, "bytes": f.getvalue()})
                if not detected_team:
                    detected_team = parsed["team"]
            else:
                unmatched_names.append(f.name)

        for k in merch_map:
            merch_map[k].sort(key=lambda x: x["name"])

        if detected_team:
            st.markdown(
                f'<div class="team-badge">'
                f'<span style="color:#888;font-size:0.7rem;letter-spacing:1px;">EQUIPO · </span>'
                f'<span style="color:#e8c84a;font-family:\'Bebas Neue\',sans-serif;letter-spacing:1px;">'
                f'{detected_team}</span></div>',
                unsafe_allow_html=True,
            )
            st.write("")

        if merch_map:
            cols = st.columns(min(4, len(merch_map)))
            for i, (k, v) in enumerate(merch_map.items()):
                with cols[i % len(cols)]:
                    st.markdown(
                        f'<div class="capsule-card match"><h4>{k}</h4><p>{len(v)} imagen(es)</p></div>',
                        unsafe_allow_html=True,
                    )

        if unmatched_names:
            with st.expander(f"⚠ {len(unmatched_names)} sin naming correcto"):
                for n in unmatched_names:
                    st.caption(n)

# ─── STEP 3: GENERAR DECK ─────────────────────────────────────
st.markdown('<div class="section-h">3 · GENERAR DECK</div>', unsafe_allow_html=True)

if not ppt_bytes:
    st.caption("Sube el PPT base en el paso 1.")
elif not merch_map:
    st.caption("Selecciona o sube los merchboards en el paso 2.")
else:
    # Preview del mapping antes de generar
    mapping_preview = []
    for s in slides_found:
        match = find_note_match(s["key"], merch_map)
        mapping_preview.append({
            "Slide":               s["slide"],
            "Nota":                s["key"],
            "Género":              s["gender"] or "—",
            "Cápsula matcheada":   match or "⚠ sin match",
            "Imágenes disponibles": len(merch_map.get(match, [])) if match else 0,
        })

    if mapping_preview:
        st.caption("Mapping propuesto:")
        st.dataframe(mapping_preview, hide_index=True, use_container_width=True)

    no_match = sum(1 for r in mapping_preview if r["Cápsula matcheada"] == "⚠ sin match")
    if no_match:
        st.warning(f"⚠ {no_match} slide(s) no tienen imágenes que coincidan. Se generará el deck pero esas slides quedarán sin cambios.")

    if st.button("Generar Deck →", type="primary"):
        with st.spinner("Procesando slides..."):
            out, replaced, log = generate_deck(ppt_bytes, merch_map)

        st.markdown(
            f'<span class="pill-ok">✓ {replaced} SLIDE(S) ACTUALIZADAS</span>',
            unsafe_allow_html=True,
        )
        st.write("")

        with st.expander("Detalle", expanded=True):
            for lvl, msg in log:
                icon  = {"ok": "✓", "warn": "⚠", "err": "✗"}.get(lvl, "·")
                color = {"ok": "#4ade80", "warn": "#e8c84a", "err": "#f87171"}.get(lvl, "#888")
                st.markdown(
                    f'<div style="font-size:0.8rem;color:{color};">{icon} {msg}</div>',
                    unsafe_allow_html=True,
                )

        base_name = re.sub(r"\.pptx$", "", ppt_file.name, flags=re.IGNORECASE)
        out_name  = f"{base_name} — {detected_team or 'EQUIPO'}.pptx"
        st.download_button(
            "⬇ Descargar PPT",
            data=out,
            file_name=out_name,
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            type="primary",
        )
